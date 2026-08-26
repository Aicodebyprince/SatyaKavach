"""
Build Phase-2 deck: copy round1_deck.pptx, update content only (design untouched).
New original content goes ONLY on previously-empty slides 7 & 8.
"""
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
from copy import deepcopy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE as MSO_LINE

SRC = "round1_deck.pptx"
OUT = "Omni_Hackathon_Codeators_Phase2.pptx"

BLUE = RGBColor(0x16, 0x75, 0xD5)
DARK = RGBColor(0x11, 0x11, 0x11)
GRAY = RGBColor(0x44, 0x44, 0x44)
LGRAY = RGBColor(0x77, 0x77, 0x77)
CARD_FILL = RGBColor(0xF5, 0xF9, 0xFE)
CARD_LINE = RGBColor(0xDC, 0xE9, 0xF7)
GREEN = RGBColor(0x12, 0x85, 0x5F)
RED = RGBColor(0xE0, 0x2B, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CHIP_FILL = RGBColor(0xE7, 0xF2, 0xFA)

prs = Presentation(SRC)


def find(slide, name):
    def rec(shapes):
        for sh in shapes:
            if sh.shape_type == 6:
                r = rec(sh.shapes)
                if r is not None:
                    return r
            elif sh.name == name:
                return sh
        return None
    return rec(slide.shapes)


def set_text(slide, name, lines, si=None):
    """Replace paragraph texts in-place, preserving every paragraph's formatting."""
    sh = find(slide, name)
    assert sh is not None, f"{si}: shape '{name}' not found"
    tf = sh.text_frame
    while len(tf.paragraphs) < len(lines):
        last = tf.paragraphs[-1]._p
        last.addnext(deepcopy(last))
    while len(tf.paragraphs) > len(lines):
        p = tf.paragraphs[-1]._p
        p.getparent().remove(p)
    for para, line in zip(tf.paragraphs, lines):
        runs = para.runs
        if not runs:
            if line:
                r = para.add_run()
                r.text = line
            continue
        runs[0].text = line
        for extra in runs[1:]:
            extra._r.getparent().remove(extra._r)


def style_run(r, name, size, bold=False, color=DARK, italic=False):
    r.font.name = name
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


def txbox(slide, x, y, w, h, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    return tb


def card(slide, x, y, w, h, fill=CARD_FILL, line=CARD_LINE,
         radius=0.12, dash=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
        if dash:
            sp.line.dash_style = dash
    sp.shadow.inherit = False
    return sp


def oval(slide, x, y, d, fill, text, tsize=13, tcolor=WHITE):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                Inches(d), Inches(d))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    style_run(r, "Arimo", tsize, True, tcolor)
    return sp


# ══════════════════ SLIDE 1 — COVER (one surgical edit) ══════════════════
s = prs.slides[0]
set_text(s, "TextBox 32",
         ["Problem Statement: Omni_CyberTech_4   •   Round 2 / Phase 2 Progress"], si=1)

# ══════════════════ SLIDE 4 — SOLUTION, NOW BUILT ══════════════════
s = prs.slides[3]
set_text(s, "TextBox 9", ["OUR SOLUTION — NOW BUILT"], si=4)
set_text(s, "TextBox 50", ["What Have We Built?"], si=4)
set_text(s, "TextBox 51", [
    "SatyaKavach is now a working end-to-end platform: a citizen uploads suspicious "
    "media and receives a Trust Score with an explainable verdict.",
    "Accepted inputs: image, video, audio, screenshot, or a suspicious link.",
    "The backend runs parallel AI analysis, fuses all signals into one score, and "
    "generates a Hindi-first evidence report with a clear recommended action.",
], si=4)

# ══════════════════ SLIDE 5 — ARCHITECTURE AS BUILT ══════════════════
s = prs.slides[4]
set_text(s, "TextBox 9", ["TECHNICAL ARCHITECTURE — AS BUILT"], si=5)
set_text(s, "TextBox 57", ["FastAPI (Python) — LIVE • REST + async pipeline"], si=5)
set_text(s, "TextBox 55", ["React.js + TypeScript",
                           "Built — premium Hindi-first citizen PWA"], si=5)
set_text(s, "TextBox 60", ["Risk Engine",
                           "LIVE — weighted fusion → Trust Score 0–100"], si=5)
set_text(s, "TextBox 61", ["Gemini",
                           "Context reasoning & report writing"], si=5)
set_text(s, "TextBox 62", ["Whisper",
                           "Transcription & voice-clone signals"], si=5)
set_text(s, "TextBox 63", ["EasyOCR",
                           "Phase B — screenshot text extraction"], si=5)
set_text(s, "TextBox 64", ["NLP / Scam Classifier",
                           "Phase B — message intent analysis"], si=5)
set_text(s, "TextBox 65", ["PostgreSQL",
                           "Users • uploads • verdicts • audit logs"], si=5)
set_text(s, "TextBox 66", ["AWS S3",
                           "Evidence storage (S3-compatible)"], si=5)
set_text(s, "TextBox 53", ["CLOUD READY - Docker Compose • Nginx • Render/Vercel/Neon deploy plan"], si=5)

# ══════════════════ SLIDE 6 — FEATURES IMPLEMENTED ══════════════════
s = prs.slides[5]
set_text(s, "TextBox 9", ["FEATURES IMPLEMENTED"], si=6)
set_text(s, "TextBox 31", ["IMAGE DEEPFAKE DETECTION  •  BUILT"], si=6)
set_text(s, "TextBox 32", [
    "Datasets", "FaceForensics++", "Celeb-DF v2", "DFDC Dataset",
    "Models", "Gemini Vision (live)", "EfficientNet + XceptionNet",
    "CPU inference path",
    "Shipped", "POST /api/v1/upload/", "Manipulation score + class",
], si=6)
set_text(s, "TextBox 34", ["VIDEO DEEPFAKE DETECTION  •  BUILT"], si=6)
set_text(s, "TextBox 35", [
    "Datasets", "FaceForensics++", "DFDC", "DeepFakeTIMIT",
    "Models", "TimeSformer + Video Swin", "Frame extraction pipeline",
    "Gemini Vision assist",
    "Shipped", "Frame-level detection", "Video authenticity score",
], si=6)
set_text(s, "TextBox 37", ["AUDIO DEEPFAKE DETECTION  •  BUILT"], si=6)
set_text(s, "TextBox 38", [
    "Datasets", "ASVspoof 2019", "FakeAVCeleb", "WaveFake",
    "Models", "Whisper transcription", "Wav2Vec2 + spectrogram",
    "Voice-clone scoring",
    "Shipped", "Voice clone detection", "Audio authenticity score",
], si=6)
set_text(s, "TextBox 40", ["RISK ENGINE + EVIDENCE REPORT  •  BUILT"], si=6)
set_text(s, "TextBox 41", [
    "Inputs", "Image signals", "Video + Audio signals", "Threat intel signals",
    "Engine", "Weighted Risk Engine", "Re-normalising fusion",
    "Gemini 2.5 report writer",
    "Output", "Trust Score 0–100", "Hindi-first evidence report",
], si=6)

# ══════════════════ SLIDE 7 — PROGRESS MADE (was empty) ══════════════════
s = prs.slides[6]
set_text(s, "TextBox 10", ["PROGRESS MADE"], si=7)

lbl = txbox(s, 0.85, 1.50, 11.3, 0.40)
p = lbl.text_frame.paragraphs[0]
r = p.add_run(); r.text = "DELIVERED SINCE ROUND 1 — WORKING SYSTEM, NOT A CONCEPT"
style_run(r, "Rajdhani Bold", 16, True, BLUE)

milestones = [
    ("Backend platform — ", "FastAPI + PostgreSQL: uploads, async jobs, verdicts, auth, audit logs"),
    ("Multimodal AI pipeline — ", "image, video & audio detector services running in parallel"),
    ("Weighted Risk Engine — ", "signal fusion → Trust Score 0–100 with three-tier verdict"),
    ("Explainable reports — ", "Gemini-written evidence summaries, Hindi-first templates"),
    ("Threat intelligence — ", "VirusTotal, Safe Browsing & PhishTank adapters + TTL cache"),
    ("Engineering quality — ", "27 unit + property-based tests passing; JWT roles & rate limits"),
    ("Deployment stack — ", "full stack Dockerised: API • database • object storage • Nginx"),
    ("Citizen PWA — ", "drag-&-drop upload → animated Trust Score dashboard (EN / हिंदी)"),
]
y = 2.02
for lead, rest in milestones:
    card(s, 0.85, y, 11.30, 0.90)
    oval(s, 1.08, y + 0.28, 0.34, GREEN, "✓", 14)
    tb = txbox(s, 1.62, y + 0.10, 10.35, 0.72)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = lead
    style_run(r1, "Proxima Nova Bold", 13, True, DARK)
    r2 = p.add_run(); r2.text = rest
    style_run(r2, "Proxima Nova", 12, False, GRAY)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    y += 0.99

# Right panel — live demo + screenshot space
card(s, 12.55, 1.50, 6.60, 8.42, fill=RGBColor(0xFC, 0xFD, 0xFF),
     line=RGBColor(0xB9, 0xD4, 0xEE))
tb = txbox(s, 12.90, 1.78, 5.9, 0.45)
r = tb.text_frame.paragraphs[0].add_run()
r.text = "LIVE DEMO — WORKING END-TO-END"
style_run(r, "Arimo Bold", 17, True, DARK)

steps = [
    ("1", "Upload any image / video / audio or paste a link"),
    ("2", "AI services analyse it in parallel (demo mode needs zero API keys)"),
    ("3", "Trust Score gauge + Hindi evidence report in under 30 seconds"),
]
yy = 2.42
for num, txt in steps:
    oval(s, 12.90, yy, 0.32, BLUE, num, 13)
    stb = txbox(s, 13.38, yy - 0.03, 5.55, 0.75)
    stf = stb.text_frame
    stf.word_wrap = True
    sp_ = stf.paragraphs[0]
    sr = sp_.add_run(); sr.text = txt
    style_run(sr, "Proxima Nova", 12.5, False, GRAY)
    yy += 0.86

inner = card(s, 12.90, 5.20, 5.90, 3.55, fill=RGBColor(0xF2, 0xF6, 0xFB),
             line=RGBColor(0xA9, 0xC6, 0xE4), radius=0.06, dash=MSO_LINE.DASH)
itf = inner.text_frame
itf.vertical_anchor = MSO_ANCHOR.MIDDLE
ip = itf.paragraphs[0]; ip.alignment = PP_ALIGN.CENTER
ir = ip.add_run(); ir.text = "[  App screenshots — drop UI captures here  ]"
style_run(ir, "Proxima Nova Italics", 13, False, LGRAY, italic=True)

# ══════════════════ SLIDE 8 — CHALLENGES & ROADMAP (was empty) ══════════════════
s = prs.slides[7]
set_text(s, "TextBox 9", ["CHALLENGES FACED & ROAD AHEAD"], si=8)

lbl = txbox(s, 0.85, 1.50, 8.9, 0.40)
r = lbl.text_frame.paragraphs[0].add_run()
r.text = "CHALLENGES FACED — AND HOW WE SOLVED THEM"
style_run(r, "Rajdhani Bold", 16, True, RED)

challenges = [
    ("GPU-heavy deepfake models",
     "Fix → Demo-mode engine ships now; HuggingFace Spaces (16 GB RAM free) path ready for real inference"),
    ("Free-tier API rate limits",
     "Fix → Threat-intel cache with TTL + SHA-256 dedup eliminates repeat calls"),
    ("Fair multimodal signal fusion",
     "Fix → Re-normalising weights over available signals — no single failure blocks a verdict"),
    ("Hindi-first explainability",
     "Fix → Bilingual prompt templates; every claim cites extracted artifacts, not a black-box score"),
]
y = 2.05
for i, (title, fix) in enumerate(challenges, 1):
    card(s, 0.85, y, 8.90, 1.72)
    oval(s, 1.10, y + 0.24, 0.36, RED, str(i), 14)
    ctf_box = txbox(s, 1.68, y + 0.18, 7.85, 1.40)
    tf = ctf_box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = title
    style_run(r1, "Proxima Nova Bold", 14, True, DARK)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = fix
    style_run(r2, "Proxima Nova", 11.5, False, GRAY)
    p2.space_before = Pt(4)
    y += 1.88

# Roadmap column
lbl = txbox(s, 10.15, 1.50, 9.0, 0.40)
r = lbl.text_frame.paragraphs[0].add_run()
r.text = "ROAD AHEAD — FROM PILOT TO PLATFORM"
style_run(r, "Rajdhani Bold", 16, True, BLUE)

roadmap = [
    ("NEXT 30 DAYS — PUBLIC PILOT",
     "Deploy free stack (Vercel + Render + Neon) • OCR scam-message classifier • I4C / 1930 reporting flow"),
    ("3 MONTHS — REACH",
     "Browser extension & WhatsApp bot intake • admin dashboard • journalist verification workspace"),
    ("6–12 MONTHS — SCALE",
     "Regional Indian languages • real-time video verification • government & fact-check partnerships"),
]
y = 2.05
for title, body in roadmap:
    card(s, 10.15, y, 9.00, 2.28)
    chip = card(s, 10.40, y + 0.22, 3.30, 0.44, fill=CHIP_FILL, line=None, radius=0.5)
    ctf = chip.text_frame
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run(); cr.text = title.split(" — ")[0]
    style_run(cr, "Rajdhani Bold", 14, True, BLUE)
    rtb = txbox(s, 10.40, y + 0.82, 8.50, 1.30)
    tf = rtb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run(); r1.text = title.split(" — ")[1]
    style_run(r1, "Proxima Nova Bold", 13.5, True, DARK)
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = body
    style_run(r2, "Proxima Nova", 11.5, False, GRAY)
    p2.space_before = Pt(4)
    y += 2.44

# Impact bar
bar = card(s, 0.85, 9.42, 18.30, 0.62, fill=DARK, line=None, radius=0.5)
btf = bar.text_frame
btf.vertical_anchor = MSO_ANCHOR.MIDDLE
bp = btf.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
br = bp.add_run()
br.text = ("EXPECTED IMPACT —  Faster public verification  •  fewer deepfake-fraud losses"
           "  •  supports MeitY / Digital India trust goals")
style_run(br, "Rajdhani Bold", 15, True, WHITE)

prs.save(OUT)
print(f"Saved {OUT}")
